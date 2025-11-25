from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle
    return slide

def add_content_slide(prs, title, content_lines):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    
    text_frame = slide.placeholders[1].text_frame
    text_frame.clear()
    
    for line in content_lines:
        p = text_frame.add_paragraph()
        if line.startswith('- '):
            p.text = line[2:]
            p.level = 0
        elif line.startswith('  - '):
            p.text = line[4:]
            p.level = 1
        else:
            p.text = line
            p.level = 0
    
    return slide

def add_table_slide(prs, title, headers, rows):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title
    
    # Add table
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(9)
    height = Inches(5)
    
    table = slide.shapes.add_table(len(rows) + 1, len(headers), left, top, width, height).table
    
    # Set column widths
    col_width = width / len(headers)
    for i in range(len(headers)):
        table.columns[i].width = int(col_width)
    
    # Add headers
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(68, 114, 196)
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.bold = True
                run.font.color.rgb = RGBColor(255, 255, 255)
                run.font.size = Pt(11)
    
    # Add rows
    for row_idx, row in enumerate(rows):
        for col_idx, cell_text in enumerate(row):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(cell_text)
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(10)
    
    return slide

# Slide 1: Title
add_title_slide(prs, 
    "AI Coding Assistants Comparison",
    "For Large Legacy Application Modernization")

# Slide 2: Executive Summary
add_content_slide(prs, "Executive Summary", [
    "- Compared five leading AI coding assistants specifically for large legacy application modernization",
    "- Key modernization factors: context handling, refactoring capabilities, code understanding",
    "- Best for monolith analysis: Cursor and Windsurf IDE (1M token context)",
    "- Best for complex refactoring: Cursor (Composer mode) and Windsurf (Cascade planning)",
    "- Best for incremental modernization: GitHub Copilot (fast code generation)",
    "- Best for secure modernization: Cline (air-gapped deployment)",
    "- Best value for large projects: Windsurf IDE (1M tokens at $15/user/month)"
])

# Slide 3: Legacy Modernization Challenges
add_content_slide(prs, "Legacy Modernization Challenges", [
    "- Understanding large codebases: Most legacy apps exceed standard context windows",
    "- Complex refactoring needs: Moving from monoliths to microservices",
    "- Technical debt: Outdated patterns, missing documentation",
    "- Cross-file dependencies: Changes affect multiple components",
    "- Maintaining business logic: Preserving critical functionality"
])

# Slide 4: Tools Evaluated
add_content_slide(prs, "Tools Evaluated for Modernization", [
    "- Claude Code: Anthropic's assistant with 200K context (Claude 3 Sonnet/Opus)",
    "- Cline: Open-source tool with self-hosted security (BYOK - any model)",
    "- GitHub Copilot: Fast code generation with 200K context (GPT-4o)",
    "- Cursor: 1M context with Composer mode for refactoring (Claude 3 Opus/GPT-4o)",
    "- Windsurf IDE: 1M context with Cascade planning for migration (Claude 3 Opus/GPT-4o)"
])

# Slide 5: Context Window Table
add_table_slide(prs, "Context Window for Legacy Codebases",
    ["Tool", "Context Window", "What It Means for Legacy Apps"],
    [
        ["Claude Code", "200K tokens", "Can analyze ~10K LOC at once"],
        ["Cline", "200K tokens", "Requires breaking large apps into modules"],
        ["GitHub Copilot", "200K tokens", "Limited understanding of full system architecture"],
        ["Cursor", "1M tokens (Max Mode)", "Can analyze ~50K LOC simultaneously"],
        ["Windsurf IDE", "1M tokens", "Can process entire monoliths at once"]
    ])

# Slide 6: Modernization Capabilities
add_table_slide(prs, "Modernization Capabilities",
    ["Tool", "Key Modernization Features", "Modernization Approach"],
    [
        ["Claude Code", "File system navigation, multi-file editing", "Conversational guidance through changes"],
        ["Cline", "Terminal access, transparency", "Task-based refactoring with control"],
        ["GitHub Copilot", "Fast code generation, pattern recognition", "Incremental modernization with quick rewrites"],
        ["Cursor", "Composer mode for multi-file refactoring", "Holistic codebase transformation"],
        ["Windsurf IDE", "Cascade planning for migration roadmaps", "Structured modernization with clear steps"]
    ])

# Slide 7: Code Understanding
add_table_slide(prs, "Legacy Code Understanding",
    ["Tool", "Code Analysis", "Dependency Mapping", "Documentation"],
    [
        ["GitHub Copilot", "Good pattern recognition", "Limited cross-file", "Good inline docs"],
        ["Cursor", "Excellent with 1M context", "Strong cross-file tracking", "Comprehensive"],
        ["Claude Code", "Good reasoning within 200K", "Manual tracking", "Excellent explanations"],
        ["Cline", "Good reasoning within 200K", "Manual tracking", "Good explanations"],
        ["Windsurf IDE", "Excellent with 1M context", "Automated mapping", "Structured docs"]
    ])

# Slide 8: Architecture & Tech Spec Generation
add_table_slide(prs, "Architecture Diagram & Tech Spec Generation",
    ["Tool", "Architecture Diagrams", "Technical Specs", "Accuracy"],
    [
        ["GitHub Copilot", "Limited (manual guidance)", "Good component-level", "Moderate"],
        ["Cursor", "Excellent with Max Mode", "Comprehensive system-wide", "Very good"],
        ["Claude Code", "Good with guidance", "Detailed but limited", "Good for documented"],
        ["Cline", "Good with guidance", "Detailed with transparency", "Good with input"],
        ["Windsurf IDE", "Excellent with Cascade", "End-to-end migration specs", "Very good"]
    ])

# Slide 9: Performance Factors
add_table_slide(prs, "Modernization Performance Factors",
    ["Tool", "Refactoring Speed", "Multi-File Changes", "Large Codebase"],
    [
        ["GitHub Copilot", "Very Fast for small changes", "Limited coordination", "Struggles with large"],
        ["Cursor", "Moderate but thorough", "Excellent with Composer", "Handles entire systems"],
        ["Windsurf IDE", "Fast with structure", "Good with Cascade", "Excellent with 1M"],
        ["Claude Code", "Moderate", "Good with manual guidance", "Limited by 200K"],
        ["Cline", "Moderate", "Good with manual control", "Limited by 200K"]
    ])

# Slide 10: Interactive Experience
add_table_slide(prs, "Interactive Modernization Experience",
    ["Tool", "Real-time Interaction", "Clarification During Processing"],
    [
        ["GitHub Copilot", "Limited (inline suggestions)", "No (requires new prompt)"],
        ["Cursor", "Excellent (continuous dialogue)", "Yes (can add clarifications)"],
        ["Windsurf IDE", "Excellent (flow mode)", "Yes (can add clarifications)"],
        ["Claude Code", "Good (conversational)", "No (must wait for completion)"],
        ["Cline", "Good (task-based)", "No (requires task completion)"]
    ])

# Slide 11: Planning Capabilities
add_table_slide(prs, "Modernization Planning Capabilities",
    ["Tool", "Plan Generation", "Plan Detail Level", "Execution Support"],
    [
        ["GitHub Copilot", "Limited (component-focused)", "Moderate", "Good code implementation"],
        ["Cursor", "Excellent (system-wide)", "Very detailed with 1M", "Excellent with guidance"],
        ["Windsurf IDE", "Excellent (structured roadmaps)", "Very detailed migration", "Excellent with Cascade"],
        ["Claude Code", "Good (conversational)", "Good but limited", "Good implementation"],
        ["Cline", "Good (transparent)", "Good with control", "Good with oversight"]
    ])

# Slide 12: Architecture Migration
add_table_slide(prs, "Legacy-to-Modern Architecture Migration",
    ["Tool", "Monolith to Microservices", "Database Modernization", "Cloud Migration"],
    [
        ["GitHub Copilot", "Good for individual services", "Good for schema updates", "Excellent for Azure"],
        ["Cursor", "Excellent system-wide", "Very good with large schemas", "Good with any cloud"],
        ["Windsurf IDE", "Excellent with roadmaps", "Very good with data modeling", "Good with any cloud"],
        ["Claude Code", "Good with guided approach", "Good within context limits", "Good with any cloud"],
        ["Cline", "Good with controlled approach", "Good within context limits", "Good with any cloud"]
    ])

# Slide 13: LLM Models
add_table_slide(prs, "LLM Models for Modernization",
    ["Tool", "Available Models", "Model Strengths"],
    [
        ["Claude Code", "Claude 3.5 Sonnet, 3 Opus, 3 Sonnet, 3 Haiku", "Excellent reasoning, strong code understanding, fast responses"],
        ["Cline", "BYOK: Claude 3.5 Sonnet, GPT-4o, o1-preview/mini, Gemini, DeepSeek, Local", "Complete flexibility, open source models, cost control"],
        ["GitHub Copilot", "GPT-4o, o1-preview/mini, Claude 3.5 Sonnet, Gemini 1.5 Pro", "Multi-model support, optimized for code completion"],
        ["Cursor", "Claude 3.5 Sonnet, Sonnet 4, GPT-4o, o1-preview/mini, Gemini 1.5 Pro/Flash", "Wide model selection, model switching, large contexts"],
        ["Windsurf IDE", "Claude 3.5 Sonnet, GPT-4o, Gemini 1.5 Pro, DeepSeek Coder, Custom", "Model switching, large contexts, cost-effective"]
    ])

# Slide 14: Enterprise Requirements
add_table_slide(prs, "Enterprise Modernization Requirements",
    ["Tool", "Legacy Code Security", "Compliance", "Team Collaboration"],
    [
        ["Claude Code", "Good data handling", "Good audit trail", "Limited sharing"],
        ["Cline", "Excellent (air-gapped)", "Excellent control", "Good with self-hosting"],
        ["GitHub Copilot", "Good with Enterprise", "Good with policy controls", "Excellent with GitHub"],
        ["Cursor", "Very good data policies", "Good audit capabilities", "Good with Business tier"],
        ["Windsurf IDE", "Good data handling", "Good tracking", "Good with Teams tier"]
    ])

# Slide 15: ROI Considerations
add_table_slide(prs, "Modernization ROI Considerations",
    ["Tool", "License Cost (10 devs/year)", "Additional Costs", "Value"],
    [
        ["Claude Code", "$3,600", "API usage limits", "Good for medium projects"],
        ["Cline", "~$2,160 (API costs)", "Self-hosting infrastructure", "Excellent for secure projects"],
        ["GitHub Copilot", "$2,280", "None significant", "Good for incremental modernization"],
        ["Cursor", "$4,800", "Max Mode credits", "Excellent for large monoliths"],
        ["Windsurf IDE", "$1,800", "None significant", "Best value for large projects"]
    ])

# Slide 16: Recommendations
add_content_slide(prs, "Legacy Modernization Recommendations", [
    "For Large Monolithic Applications (>300K LOC):",
    "  - Best Choice: Cursor or Windsurf IDE",
    "  - 1M token context essential for understanding entire system",
    "",
    "For Incremental Modernization:",
    "  - Best Choice: GitHub Copilot",
    "  - Faster for smaller, targeted improvements",
    "",
    "For Regulated Industries:",
    "  - Best Choice: Cline",
    "  - Air-gapped deployment for sensitive codebases",
    "",
    "For Cost-Effective Large Projects:",
    "  - Best Choice: Windsurf IDE",
    "  - Lowest cost with 1M token capability"
])

# Slide 17: Implementation Plan
add_content_slide(prs, "Legacy Modernization Implementation Plan", [
    "1. Assessment Phase (2-4 weeks)",
    "  - Analyze legacy codebase size and complexity",
    "  - Identify modernization goals",
    "  - Select appropriate AI tool",
    "",
    "2. Pilot Modernization (4-8 weeks)",
    "  - Select isolated component for initial modernization",
    "  - Use AI tool to analyze, document, and refactor",
    "",
    "3. Scaled Modernization (3-6 months)",
    "  - Develop modernization patterns with AI assistance",
    "  - Create reusable templates",
    "",
    "4. Continuous Modernization",
    "  - Integrate AI tools into ongoing development",
    "  - Measure business impact"
])

# Slide 18: Tool Selection Matrix
add_table_slide(prs, "Tool Selection Matrix",
    ["Modernization Scenario", "Best Tool", "Second Choice"],
    [
        ["Large monolith (>300K LOC)", "Cursor", "Windsurf IDE"],
        ["Security-critical systems", "Cline", "GitHub Copilot Enterprise"],
        ["Rapid incremental updates", "GitHub Copilot", "Claude Code"],
        ["Microservices migration", "Windsurf IDE", "Cursor"],
        ["Limited modernization budget", "Windsurf IDE", "GitHub Copilot"],
        ["Database modernization", "Cursor", "Windsurf IDE"],
        ["Cloud migration (Azure)", "GitHub Copilot", "Cursor"],
        ["Legacy code documentation", "Claude Code", "Cursor"]
    ])

# Save presentation
prs.save('AI-Coding-Assistants-Presentation.pptx')
print("PowerPoint presentation created successfully: AI-Coding-Assistants-Presentation.pptx")

